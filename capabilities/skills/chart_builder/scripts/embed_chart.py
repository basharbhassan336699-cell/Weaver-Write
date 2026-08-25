"""
embed_chart.py — render a chart AND embed it into a document (working script)
=============================================================================
Closes the gap between the `chart` builder and the document builders:
generate a themed chart image, then insert it directly into a Word (.docx)
or PowerPoint (.pptx) file at the right place — captioned and centered.

This is what makes "put a chart of these results into my report" one step.

Functions:
    embed_chart_in_docx(docx_path, chart_spec, caption, ...) -> updates the docx
    add_chart_slide(pptx_path, chart_spec, title, ...)       -> adds a chart slide
    make_chart_image(chart_spec, out_png)                    -> just the image

Requires: matplotlib, python-docx, python-pptx
"""
from __future__ import annotations
import os
import tempfile

_SCRIPTS = os.path.dirname(os.path.abspath(__file__))
import sys
if _SCRIPTS not in sys.path:
    sys.path.insert(0, _SCRIPTS)


def make_chart_image(chart_spec, out_png=None):
    """chart_spec: {type, data, title?, theme_id?, xlabel?, ylabel?, lang?}."""
    from build_chart import build_chart
    if out_png is None:
        fd, out_png = tempfile.mkstemp(suffix=".png")
        os.close(fd)
    r = build_chart(
        chart_spec.get("type", "bar"), chart_spec.get("data", {}), out_png,
        title=chart_spec.get("title", ""),
        theme_id=chart_spec.get("theme_id", "academic_navy"),
        xlabel=chart_spec.get("xlabel", ""), ylabel=chart_spec.get("ylabel", ""),
        lang=chart_spec.get("lang", "ar"),
        dpi=chart_spec.get("dpi", 150))
    return r


def embed_chart_in_docx(docx_path, chart_spec, caption="", lang="ar",
                        width_inches=6.0, output_path=None):
    """
    Render a chart and append it (centered) with a caption to an existing docx.
    If the docx doesn't exist, a new one is created.
    """
    from docx import Document
    from docx.shared import Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    r = make_chart_image(chart_spec)
    if not r.get("ok"):
        return r
    img = r["output_path"]

    doc = Document(docx_path) if os.path.exists(docx_path) else Document()
    # image, centered
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run()
    run.add_picture(img, width=Inches(width_inches))
    # caption
    if caption:
        cap = doc.add_paragraph(caption)
        cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
        for r_ in cap.runs:
            r_.font.italic = True
            r_.font.size = __import__("docx").shared.Pt(10)

    out = output_path or docx_path
    doc.save(out)
    try: os.unlink(img)
    except Exception: pass
    return {"ok": True, "output_path": out, "embedded": chart_spec.get("type"),
            "engine": "python-docx"}


def add_chart_slide(pptx_path, chart_spec, title="", lang="ar", output_path=None):
    """
    Render a chart and add it as a new slide (image centered) to a pptx.
    Creates the pptx if it doesn't exist.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    r = make_chart_image(chart_spec)
    if not r.get("ok"):
        return r
    img = r["output_path"]

    prs = Presentation(pptx_path) if os.path.exists(pptx_path) else Presentation()
    if not os.path.exists(pptx_path):
        prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    # pick a blank layout; fall back to the last available layout if index 6
    # doesn't exist (html2pptx decks may have few layouts)
    layouts = prs.slide_layouts
    try:
        blank = layouts[6]
    except (IndexError, KeyError):
        blank = layouts[len(layouts) - 1]
    slide = prs.slides.add_slide(blank)

    # title
    if title:
        tb = slide.shapes.add_textbox(Inches(0.7), Inches(0.4),
                                      prs.slide_width - Inches(1.4), Inches(0.9))
        tf = tb.text_frame; tf.word_wrap = True
        para = tf.paragraphs[0]; run = para.add_run(); run.text = title
        run.font.size = Pt(26); run.font.bold = True
        para.alignment = PP_ALIGN.RIGHT if lang == "ar" else PP_ALIGN.LEFT

    # center the chart image
    from PIL import Image
    iw, ih = Image.open(img).size
    ratio = ih / iw
    disp_w = Inches(9)
    disp_h = Inches(9 * ratio)
    left = (prs.slide_width - disp_w) / 2
    top = Inches(1.5)
    slide.shapes.add_picture(img, left, top, width=disp_w, height=disp_h)

    out = output_path or pptx_path
    prs.save(out)
    try: os.unlink(img)
    except Exception: pass
    return {"ok": True, "output_path": out, "embedded": chart_spec.get("type"),
            "engine": "python-pptx"}


if __name__ == "__main__":
    import argparse, json
    p = argparse.ArgumentParser(description="Embed a chart into a document")
    p.add_argument("--target", required=True, help="docx or pptx path")
    p.add_argument("--spec", required=True, help="JSON chart spec")
    p.add_argument("--caption", default="")
    p.add_argument("--lang", default="ar")
    args = p.parse_args()
    with open(args.spec, encoding="utf-8") as f:
        spec = json.load(f)
    if args.target.endswith(".pptx"):
        r = add_chart_slide(args.target, spec, title=args.caption, lang=args.lang)
    else:
        r = embed_chart_in_docx(args.target, spec, caption=args.caption, lang=args.lang)
    print(json.dumps(r, ensure_ascii=False))
