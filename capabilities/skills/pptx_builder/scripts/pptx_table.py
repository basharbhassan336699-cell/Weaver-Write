"""
pptx_table.py — native PowerPoint tables with correct direction (working)
=========================================================================
Adds a REAL PowerPoint table (a:tbl) to a slide, themed and direction-aware:

  - Arabic (lang="ar")  -> RTL: the whole table is marked rtl="1", columns
    are emitted right-to-left (first data column on the RIGHT), cell text is
    right-aligned. This is what makes an Arabic table "start from the right".
  - English (lang="en") -> LTR: normal left-to-right.

Themed header row (primary fill, white bold), optional totals row (accent).

Works on an existing pptx (adds a slide) or builds a one-slide file.

Requires: python-pptx
"""
from __future__ import annotations
import os
import json

_THEMES = os.path.join(
    os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))),
    "themes", "themes.json")


def _palette(theme_id="academic_navy"):
    try:
        with open(_THEMES, encoding="utf-8") as f:
            t = json.load(f)["themes"].get(theme_id, {})
        return (t.get("primary", "1B2A4A"), t.get("accent", "C8A04A"),
                t.get("text", "222A38"))
    except Exception:
        return ("1B2A4A", "C8A04A", "222A38")


def _rgb(hex_str):
    from pptx.dml.color import RGBColor
    return RGBColor.from_string(hex_str)


def _set_cell_rtl(cell, rtl):
    """Set paragraph direction inside a table cell."""
    from pptx.oxml.ns import qn
    tc = cell._tc
    txBody = tc.find(qn("a:txBody"))
    if txBody is None:
        return
    for p in txBody.findall(qn("a:p")):
        pPr = p.find(qn("a:pPr"))
        if pPr is None:
            pPr = p.makeelement(qn("a:pPr"), {})
            p.insert(0, pPr)
        pPr.set("rtl", "1" if rtl else "0")
        pPr.set("algn", "r" if rtl else "l")


def add_table_slide(prs_or_path, headers, rows, lang="ar",
                    theme_id="academic_navy", title="", totals=None,
                    font=None, output_path=None):
    """
    Add a slide with a native, direction-correct table.
    Returns dict with ok/output_path.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn

    rtl = (lang == "ar")
    primary, accent, text_col = _palette(theme_id)
    font = font or ("Kufyan Arabic" if rtl else "Georgia")

    # open or create
    if isinstance(prs_or_path, str) and os.path.exists(prs_or_path):
        prs = Presentation(prs_or_path)
    elif isinstance(prs_or_path, str):
        prs = Presentation()
        prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    else:
        prs = prs_or_path

    layouts = prs.slide_layouts
    blank = layouts[6] if len(layouts) > 6 else layouts[-1]
    slide = prs.slides.add_slide(blank)

    # RTL: reverse column order so the first logical column sits on the RIGHT
    display_headers = list(headers)
    display_rows = [list(r) for r in rows]
    display_totals = list(totals) if totals else None
    if rtl:
        display_headers = display_headers[::-1]
        display_rows = [r[::-1] for r in display_rows]
        if display_totals:
            display_totals = display_totals[::-1]

    # title
    if title:
        tb = slide.shapes.add_textbox(Inches(0.7), Inches(0.4),
                                      prs.slide_width - Inches(1.4), Inches(0.8))
        tf = tb.text_frame; tf.word_wrap = True
        para = tf.paragraphs[0]; run = para.add_run(); run.text = title
        run.font.size = Pt(26); run.font.bold = True
        run.font.name = font; run.font.color.rgb = _rgb(primary)
        para.alignment = PP_ALIGN.RIGHT if rtl else PP_ALIGN.LEFT

    n_rows = 1 + len(display_rows) + (1 if display_totals else 0)
    n_cols = len(display_headers)
    gfx = slide.shapes.add_table(n_rows, n_cols, Inches(0.7), Inches(1.5),
                                 prs.slide_width - Inches(1.4), Inches(0.5 * n_rows))
    table = gfx.table

    # mark the whole table RTL at XML level
    tblPr = table._tbl.find(qn("a:tblPr"))
    if tblPr is not None:
        tblPr.set("rtl", "1" if rtl else "0")

    # header row
    for c, h in enumerate(display_headers):
        cell = table.cell(0, c)
        cell.fill.solid(); cell.fill.fore_color.rgb = _rgb(primary)
        cell.text = str(h)
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.bold = True; run.font.size = Pt(13)
                run.font.color.rgb = _rgb("FFFFFF"); run.font.name = font
        _set_cell_rtl(cell, rtl)

    # body rows
    for r, row in enumerate(display_rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.RIGHT if rtl else PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.size = Pt(12); run.font.name = font
                    run.font.color.rgb = _rgb(text_col)
            _set_cell_rtl(cell, rtl)

    # totals row
    if display_totals:
        r = n_rows - 1
        for c, val in enumerate(display_totals):
            cell = table.cell(r, c)
            cell.fill.solid(); cell.fill.fore_color.rgb = _rgb(accent)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.RIGHT if rtl else PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.bold = True; run.font.size = Pt(12)
                    run.font.color.rgb = _rgb("1A1A1A"); run.font.name = font
            _set_cell_rtl(cell, rtl)

    if isinstance(prs_or_path, str):
        out = output_path or prs_or_path
        os.makedirs(os.path.dirname(out) or ".", exist_ok=True)
        prs.save(out)
        return {"ok": True, "output_path": out, "direction": "RTL" if rtl else "LTR",
                "engine": "python-pptx-table"}
    return {"ok": True, "slide": slide, "direction": "RTL" if rtl else "LTR"}


if __name__ == "__main__":
    r = add_table_slide("/tmp/native_table.pptx",
                        ["البند", "القيمة", "النسبة"],
                        [["الأول", 100, "25%"], ["الثاني", 200, "50%"]],
                        lang="ar", title="جدول البيانات",
                        totals=["الإجمالي", 300, "75%"])
    print(r)
