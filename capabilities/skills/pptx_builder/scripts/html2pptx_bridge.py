"""
html2pptx_bridge.py — convert themed HTML to native PPTX (working script)
=========================================================================
Bridges Python -> the vendored Node.js html2pptx engine
(engines/html2pptx-core). Given an HTML deck, it produces a native,
editable .pptx (not an image), preserving the CSS design.

This is the "Claude way": design in HTML/CSS, convert to real PPTX shapes.

Requires: Node.js + the engine deps (adm-zip, cheerio, css, pptxgenjs),
installed once in engines/html2pptx-core/.
"""
from __future__ import annotations
import argparse
import os
import re
import subprocess
import tempfile
import shutil
from copy import deepcopy

_ENGINE = os.path.join(
    os.path.dirname(os.path.dirname(os.path.dirname(os.path.dirname(
        os.path.dirname(os.path.abspath(__file__)))))),
    "engines", "html2pptx-core")


def html_to_pptx(html_path_or_str, output_path, is_string=False):
    """
    Convert an HTML file (or HTML string) to native PPTX.

    Backward-compatible entry point. A file input, or a single-slide string,
    behaves EXACTLY as before (one Node conversion). A multi-slide HTML string
    (several `<div class="slide">`) is split, each slide converted via the Node
    engine, and the parts merged — working around the engine's one-page limit
    without touching it. Returns
    {"ok": bool, "output_path": str, "engine": str, "slides": int?, "error": str?}
    """
    if is_string:
        return html_to_pptx_multi(html_path_or_str, output_path,
                                  _html_to_pptx_single)
    return _html_to_pptx_single(html_path_or_str, output_path, is_string=False)


def split_slides_html(html: str):
    """Return a list of HTML documents, each holding ONE slide + the same
    <head>/<style>. If there is only one slide, return [html] unchanged (so the
    original single-conversion path is used). (tested)"""
    head_m = re.search(r"<head>.*?</head>", html, re.S | re.I)
    head = head_m.group(0) if head_m else "<head></head>"
    slides = re.findall(
        r'<div class="slide".*?</div>\s*(?=<div class="slide"|</body>|</html>|$)',
        html, re.S | re.I)
    if len(slides) <= 1:
        return [html]
    return [f"<!DOCTYPE html><html>{head}<body>{s}</body></html>"
            for s in slides]


def html_to_pptx_multi(html_str: str, output_path: str, single_fn):
    """Convert HTML (possibly several slides) into ONE PPTX holding them all.

    single_fn: the original single-conversion function
    html_to_pptx(html, out, is_string=True). Returns the same result shape as
    the bridge: {"ok":bool,"output_path":str,"engine":str,"slides":int?,"error":str?}
    (tested: 1 slide -> 1, 5 slides -> 5)."""
    docs = split_slides_html(html_str)

    # single slide or un-splittable -> original path (full backward compat)
    if len(docs) <= 1:
        return single_fn(html_str, output_path, is_string=True)

    from pptx import Presentation

    parts, tmps = [], []
    try:
        for i, doc in enumerate(docs):
            fd, tmp = tempfile.mkstemp(suffix=f"_part{i}.pptx")
            os.close(fd)
            tmps.append(tmp)
            r = single_fn(doc, tmp, is_string=True)
            if r.get("ok") and os.path.exists(tmp):
                parts.append(tmp)

        if not parts:
            return {"ok": False, "error": "no slides converted"}

        merged = Presentation(parts[0])           # first part = the base
        for p in parts[1:]:                        # append the rest's slides
            src = Presentation(p)
            blank = (merged.slide_layouts[6]
                     if len(merged.slide_layouts) > 6
                     else merged.slide_layouts[-1])
            for slide in src.slides:
                new_slide = merged.slides.add_slide(blank)
                for shape in slide.shapes:
                    new_slide.shapes._spTree.append(deepcopy(shape._element))

        output_path = os.path.abspath(output_path)
        os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
        merged.save(output_path)
        return {"ok": True, "output_path": output_path,
                "engine": "html2pptx", "slides": len(parts)}
    finally:
        for t in tmps:
            try:
                os.unlink(t)
            except Exception:
                pass


def _html_to_pptx_single(html_path_or_str, output_path, is_string=False):
    """
    Convert ONE HTML file (or single-slide HTML string) to native PPTX via the
    Node engine. This is the original html_to_pptx logic, unchanged.
    Returns {"ok": bool, "output_path": str, "engine": str, "error": str?}
    """
    node = shutil.which("node")
    if not node:
        return {"ok": False, "error": "Node.js not installed (needed for html2pptx). "
                                      "Termux: pkg install nodejs"}

    lib = os.path.join(_ENGINE, "lib", "html2pptx.js")
    if not os.path.exists(lib):
        return {"ok": False, "error": f"html2pptx engine not found at {lib}"}

    # engine deps must be installed
    if not os.path.isdir(os.path.join(_ENGINE, "node_modules")):
        return {"ok": False,
                "error": "html2pptx deps not installed. Run once: "
                         f"cd {_ENGINE} && npm install"}

    # materialize HTML to a temp file if given as string
    cleanup = None
    if is_string:
        fd, html_path = tempfile.mkstemp(suffix=".html")
        with os.fdopen(fd, "w", encoding="utf-8") as f:
            f.write(html_path_or_str)
        cleanup = html_path
    else:
        html_path = html_path_or_str

    output_path = os.path.abspath(output_path)
    os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)

    # small Node runner
    runner = (
        "const {convertHTML2PPTX}=require(process.argv[1]);"
        "convertHTML2PPTX(process.argv[2],process.argv[3])"
        ".then(r=>{console.log('OK');process.exit(0);})"
        ".catch(e=>{console.error(e.message);process.exit(1);});"
    )
    try:
        proc = subprocess.run(
            [node, "-e", runner, lib, html_path, output_path],
            capture_output=True, text=True, timeout=120, cwd=_ENGINE)
    finally:
        if cleanup:
            try: os.unlink(cleanup)
            except Exception: pass

    if proc.returncode != 0:
        return {"ok": False, "error": proc.stderr.strip() or "conversion failed"}
    return {"ok": True, "output_path": output_path, "engine": "html2pptx"}


def _main():
    p = argparse.ArgumentParser(description="Convert HTML to native PPTX")
    p.add_argument("--input", required=True, help="HTML file path")
    p.add_argument("--output", default="deck.pptx")
    args = p.parse_args()
    r = html_to_pptx(args.input, args.output)
    if r["ok"]:
        print(f"Created: {r['output_path']}")
    else:
        print(f"Error: {r['error']}")


if __name__ == "__main__":
    _main()
