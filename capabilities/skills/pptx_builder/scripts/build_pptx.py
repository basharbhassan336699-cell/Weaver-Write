"""
build_pptx.py — build a professional PowerPoint deck (working script)
=====================================================================
High-fidelity slide builder with correct bidirectional handling:
  - Arabic (lang="ar")  -> RTL: right-aligned, paragraph direction RTL.
  - English (lang="en") -> LTR: left-aligned, left-to-right.
Direction is set from `lang` and auto-detected per run for mixed content.

Design system (navy/gold academic theme): title / section / content /
closing slides, consistent palette and typography.

Requires: pip install python-pptx
"""
from __future__ import annotations
import argparse
import json
import re

NAVY  = (0x1B, 0x2A, 0x4A)
GOLD  = (0xC8, 0xA0, 0x4A)
LIGHT = (0xF4, 0xF6, 0xFA)
DARK  = (0x22, 0x2A, 0x38)
WHITE = (0xFF, 0xFF, 0xFF)

AR_FONT = "Kufyan Arabic"
EN_FONT = "Georgia"

_AR_RANGE = re.compile(r"[\u0600-\u06FF\u0750-\u077F\u08A0-\u08FF\uFB50-\uFDFF\uFE70-\uFEFF]")


def _is_rtl_text(text: str) -> bool:
    if not text:
        return False
    ar = len(_AR_RANGE.findall(text))
    letters = len([c for c in text if c.isalpha()])
    return letters > 0 and (ar / max(letters, 1)) >= 0.4


def _rgb(t):
    from pptx.dml.color import RGBColor
    return RGBColor(*t)


def _set_paragraph_rtl(paragraph, rtl: bool):
    pPr = paragraph._p.get_or_add_pPr()
    pPr.set("rtl", "1" if rtl else "0")
    return pPr


def _apply_run_style(run, *, size, bold, color, font_ar, font_en, rtl):
    from pptx.util import Pt
    from pptx.oxml.ns import qn
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)
    fam = font_ar if (_is_rtl_text(run.text) or rtl) else font_en
    run.font.name = fam
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", fam)


def _add_text(frame, text, *, size, bold, color, align, rtl,
              font_ar=AR_FONT, font_en=EN_FONT, clear=True):
    from pptx.enum.text import PP_ALIGN
    if clear and frame.paragraphs and not frame.paragraphs[0].runs:
        p = frame.paragraphs[0]
    else:
        p = frame.add_paragraph()
    run = p.add_run()
    run.text = text
    _apply_run_style(run, size=size, bold=bold, color=color,
                     font_ar=font_ar, font_en=font_en, rtl=rtl)
    if align == "center":
        p.alignment = PP_ALIGN.CENTER
    elif align == "right" or (align is None and rtl):
        p.alignment = PP_ALIGN.RIGHT
    else:
        p.alignment = PP_ALIGN.LEFT
    _set_paragraph_rtl(p, rtl)
    return p


def _bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(color)


def _accent_bar(slide, prs, *, rtl, color=GOLD, width_in=2.2, height_in=0.12, top_in=1.55):
    from pptx.util import Inches
    from pptx.enum.shapes import MSO_SHAPE
    slide_w = prs.slide_width
    left = (slide_w - Inches(0.9) - Inches(width_in)) if rtl else Inches(0.9)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(top_in),
                                   Inches(width_in), Inches(height_in))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(color)
    shape.line.fill.background()
    return shape


def _title_slide(prs, title, subtitle, rtl, fa, fe):
    from pptx.util import Inches
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, NAVY)
    align = "right" if rtl else "left"
    box = slide.shapes.add_textbox(Inches(0.9), Inches(2.4),
                                   prs.slide_width - Inches(1.8), Inches(1.6))
    box.text_frame.word_wrap = True
    _add_text(box.text_frame, title, size=40, bold=True, color=WHITE,
              align=align, rtl=rtl, font_ar=fa, font_en=fe)
    if subtitle:
        sb = slide.shapes.add_textbox(Inches(0.9), Inches(4.0),
                                      prs.slide_width - Inches(1.8), Inches(1.0))
        sb.text_frame.word_wrap = True
        _add_text(sb.text_frame, subtitle, size=20, bold=False, color=GOLD,
                  align=align, rtl=rtl, font_ar=fa, font_en=fe)
    _accent_bar(slide, prs, rtl=rtl, top_in=2.2)
    return slide


def _section_slide(prs, title, rtl, fa, fe):
    from pptx.util import Inches
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, GOLD)
    align = "right" if rtl else "left"
    box = slide.shapes.add_textbox(Inches(0.9), Inches(3.0),
                                   prs.slide_width - Inches(1.8), Inches(1.4))
    box.text_frame.word_wrap = True
    _add_text(box.text_frame, title, size=34, bold=True, color=NAVY,
              align=align, rtl=rtl, font_ar=fa, font_en=fe)
    return slide


def _content_slide(prs, title, points, rtl, fa, fe):
    from pptx.util import Inches
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, LIGHT)
    align = "right" if rtl else "left"
    tb = slide.shapes.add_textbox(Inches(0.9), Inches(0.6),
                                  prs.slide_width - Inches(1.8), Inches(1.0))
    tb.text_frame.word_wrap = True
    _add_text(tb.text_frame, title, size=28, bold=True, color=NAVY,
              align=align, rtl=rtl, font_ar=fa, font_en=fe)
    _accent_bar(slide, prs, rtl=rtl, top_in=1.45, width_in=1.6)
    body = slide.shapes.add_textbox(Inches(0.9), Inches(1.9),
                                    prs.slide_width - Inches(1.8),
                                    prs.slide_height - Inches(2.6))
    tf = body.text_frame
    tf.word_wrap = True
    for i, pt in enumerate(points):
        p_rtl = True if rtl else _is_rtl_text(pt)
        marker = "◀ " if p_rtl else "▶ "
        _add_text(tf, marker + pt, size=18, bold=False, color=DARK,
                  align=("right" if p_rtl else "left"), rtl=p_rtl,
                  font_ar=fa, font_en=fe, clear=(i == 0))
    return slide


def _closing_slide(prs, text, rtl, fa, fe):
    from pptx.util import Inches
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, NAVY)
    box = slide.shapes.add_textbox(Inches(0.9), Inches(3.1),
                                   prs.slide_width - Inches(1.8), Inches(1.2))
    box.text_frame.word_wrap = True
    _add_text(box.text_frame, text, size=32, bold=True, color=GOLD,
              align="center", rtl=rtl, font_ar=fa, font_en=fe)
    return slide


def build_deck(title, slides, subtitle="", output_path="deck.pptx",
               lang="ar", closing=None):
    """Build a complete professional deck. lang: 'ar' (RTL) | 'en' (LTR)."""
    from pptx import Presentation
    from pptx.util import Inches
    rtl = (lang == "ar")
    fa, fe = AR_FONT, EN_FONT
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    _title_slide(prs, title, subtitle, rtl, fa, fe)
    for s in slides:
        if s.get("layout") == "section":
            _section_slide(prs, s.get("title", ""), rtl, fa, fe)
        else:
            _content_slide(prs, s.get("title", ""), s.get("points", []), rtl, fa, fe)
    if closing is None:
        closing = "شكراً لكم" if rtl else "Thank you"
    _closing_slide(prs, closing, rtl, fa, fe)
    prs.save(output_path)
    return output_path


def build_pptx(slides, output_path, lang="ar", title="", subtitle="", closing=None):
    """Alias used by the doc_export tool."""
    return build_deck(title=title, slides=slides, subtitle=subtitle,
                      output_path=output_path, lang=lang, closing=closing)


def _main():
    p = argparse.ArgumentParser(description="Build a professional PowerPoint deck")
    p.add_argument("--json", required=True)
    p.add_argument("--output", default="deck.pptx")
    p.add_argument("--lang", default="ar", choices=["ar", "en"])
    args = p.parse_args()
    with open(args.json, encoding="utf-8") as f:
        data = json.load(f)
    path = build_deck(title=data.get("title", ""), subtitle=data.get("subtitle", ""),
                      slides=data.get("slides", []), output_path=args.output,
                      lang=args.lang, closing=data.get("closing"))
    print(f"Created: {path}")


if __name__ == "__main__":
    _main()
